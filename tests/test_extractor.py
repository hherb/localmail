# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Horst Herb

"""Tests for the extractor protocol and LightweightExtractor."""

from __future__ import annotations

import dataclasses
import io
import logging
from pathlib import Path

import pytest

from localmail.search.extractor import (
    AttachmentExtractor,
    ExtractedText,
    ExtractorError,
    LightweightExtractor,
)


def test_extracted_text_is_frozen_dataclass() -> None:
    et = ExtractedText(text="hello", page_count=1, extractor="x@1")
    with pytest.raises(dataclasses.FrozenInstanceError):
        et.text = "world"  # type: ignore[misc]


class TestExtractedTextIsDbSafeByConstruction:
    """``ExtractedText.text`` never carries a NUL byte (#249).

    Postgres TEXT rejects ``\\x00``, and ``extracted_text`` is the only
    consumer. Enforcing it here rather than in each of the eleven
    ``_extract_*`` methods means a twelfth cannot forget: a NUL that reached
    the INSERT aborted it, escaped to the worker's outer safety net, and was
    recorded as a poison pill under the extractor name ``'unexpected'`` —
    permanently, since re-extracting the same bytes reproduces the same NUL.
    """

    def test_a_nul_in_the_text_is_stripped(self) -> None:
        assert ExtractedText(
            text="before\x00after", page_count=None, extractor="x@1"
        ).text == "beforeafter"

    def test_every_nul_is_stripped(self) -> None:
        assert ExtractedText(
            text="\x00a\x00b\x00", page_count=None, extractor="x@1"
        ).text == "ab"

    def test_clean_text_is_untouched(self) -> None:
        assert ExtractedText(
            text="ordinary text", page_count=3, extractor="x@1"
        ).text == "ordinary text"

    def test_text_of_only_nuls_becomes_the_empty_sentinel(self) -> None:
        """'' already means "we tried, got nothing, don't retry" — the right
        outcome for a blob whose entire extracted text was NUL bytes."""
        assert ExtractedText(
            text="\x00\x00", page_count=None, extractor="x@1"
        ).text == ""

    def test_the_dataclass_stays_frozen(self) -> None:
        """Normalising in __post_init__ must not have unfrozen the class."""
        et = ExtractedText(text="a\x00b", page_count=None, extractor="x@1")
        with pytest.raises(dataclasses.FrozenInstanceError):
            et.text = "c"  # type: ignore[misc]

    def test_equality_still_works_after_normalising(self) -> None:
        assert ExtractedText(
            text="a\x00b", page_count=1, extractor="x@1"
        ) == ExtractedText(text="ab", page_count=1, extractor="x@1")


class TestWhitespaceOnlyTextBecomesTheSentinel:
    """Whitespace-only text normalises to '' on construction (#266).

    The attachment-chunking claim skips sentinel rows via
    ``extracted_text <> ''``, but ``chunk_attachment_text`` yields no chunks
    for any text whose whitespace-normalisation is empty — so a stored
    whitespace-only row passes the claim, produces nothing, and is re-claimed
    on every sweep forever. Enough of them sorting low in the sha256 order
    fills the batch and attachment ingestion silently stops (the #216 shape).
    Same by-construction placement as the NUL rule above.
    """

    def test_spaces_tabs_and_newlines_become_empty(self) -> None:
        assert ExtractedText(
            text=" \t\n \n\t ", page_count=None, extractor="x@1"
        ).text == ""

    def test_unicode_whitespace_becomes_empty(self) -> None:
        """A page of non-breaking spaces is the observed real-world case."""
        assert ExtractedText(
            text="\u00a0\u00a0\n\u00a0\u2003", page_count=None, extractor="x@1"
        ).text == ""

    def test_text_with_any_substance_keeps_its_whitespace(self) -> None:
        """Only fully-blank text is collapsed; real text is stored verbatim
        (chunking does its own normalisation later)."""
        assert ExtractedText(
            text="  x  \n", page_count=None, extractor="x@1"
        ).text == "  x  \n"

    def test_nuls_amid_whitespace_still_become_empty(self) -> None:
        """The NUL strip runs first; what remains is whitespace-only."""
        assert ExtractedText(
            text=" \x00 \n\x00", page_count=None, extractor="x@1"
        ).text == ""


def test_lightweight_supports_pdf_mime_and_ext() -> None:
    lw = LightweightExtractor()
    assert lw.supports("application/pdf", "foo.pdf")
    assert lw.supports(None, "foo.pdf")
    assert lw.supports("application/pdf", "")
    assert not lw.supports("video/mp4", "foo.mp4")
    assert not lw.supports(None, None)


def test_lightweight_does_not_support_image() -> None:
    lw = LightweightExtractor()
    assert not lw.supports("image/png", "logo.png")


def test_extractor_protocol_runtime_checkable() -> None:
    """AttachmentExtractor is @runtime_checkable so workers can verify
    a configured class implements the interface."""
    lw = LightweightExtractor()
    assert isinstance(lw, AttachmentExtractor)


def _build_native_pdf(text: str) -> bytes:
    """Build a single-page text PDF in memory using reportlab."""
    from reportlab.pdfgen import canvas
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(72, 720, text)
    c.showPage()
    c.save()
    return buf.getvalue()


def test_lightweight_extracts_native_pdf(tmp_path) -> None:
    pdf_bytes = _build_native_pdf("hello attachment world")
    p = tmp_path / "a.pdf"
    p.write_bytes(pdf_bytes)

    lw = LightweightExtractor()
    result = lw.extract(p, "application/pdf")

    assert "hello attachment world" in result.text
    assert result.extractor == "lightweight@1.0"
    assert result.page_count == 1


def test_lightweight_returns_empty_on_scanned_pdf(tmp_path) -> None:
    """A PDF whose only content is a rasterized image of text returns ''
    from pypdf — the docling fallback exists for this case."""
    from PIL import Image, ImageDraw
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    img = Image.new("RGB", (400, 80), "white")
    ImageDraw.Draw(img).text((10, 30), "scanned text content", fill="black")
    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_buf.seek(0)

    pdf_buf = io.BytesIO()
    c = canvas.Canvas(pdf_buf)
    c.drawImage(ImageReader(img_buf), 72, 600, width=400, height=80)
    c.showPage()
    c.save()

    p = tmp_path / "scan.pdf"
    p.write_bytes(pdf_buf.getvalue())

    lw = LightweightExtractor()
    result = lw.extract(p, "application/pdf")

    assert result.text == ""


def test_lightweight_raises_on_encrypted_pdf(tmp_path) -> None:
    import pikepdf
    pdf_bytes = _build_native_pdf("secret")
    src = tmp_path / "src.pdf"
    src.write_bytes(pdf_bytes)

    enc = tmp_path / "enc.pdf"
    with pikepdf.open(src) as p:
        p.save(enc, encryption=pikepdf.Encryption(owner="o", user="u", R=4))

    lw = LightweightExtractor()
    with pytest.raises(ExtractorError):
        lw.extract(enc, "application/pdf")


def test_lightweight_extracts_docx(tmp_path) -> None:
    import docx
    p = tmp_path / "a.docx"
    d = docx.Document()
    d.add_paragraph("docx paragraph one")
    d.add_paragraph("docx paragraph two")
    d.save(str(p))

    lw = LightweightExtractor()
    result = lw.extract(
        p,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert "docx paragraph one" in result.text
    assert "docx paragraph two" in result.text
    assert result.extractor == "lightweight@1.0"


def test_lightweight_extracts_xlsx(tmp_path) -> None:
    import openpyxl
    p = tmp_path / "a.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "row one cell A"
    ws["B1"] = "row one cell B"
    ws2 = wb.create_sheet(title="Sheet2")
    ws2["A1"] = "second sheet content"
    wb.save(str(p))

    lw = LightweightExtractor()
    result = lw.extract(
        p,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    assert "row one cell A" in result.text
    assert "row one cell B" in result.text
    assert "second sheet content" in result.text
    assert result.extractor == "lightweight@1.0"


def test_lightweight_extracts_pptx(tmp_path) -> None:
    from pptx import Presentation
    p = tmp_path / "a.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = "Slide title here"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "speaker note content"
    prs.save(str(p))

    lw = LightweightExtractor()
    result = lw.extract(
        p,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    assert "Slide title here" in result.text
    assert "speaker note content" in result.text


def test_lightweight_extracts_txt_utf8(tmp_path) -> None:
    p = tmp_path / "a.txt"
    p.write_text("plain text content über alles", encoding="utf-8")

    lw = LightweightExtractor()
    result = lw.extract(p, "text/plain")
    assert "plain text content über alles" == result.text


def test_lightweight_extracts_txt_latin1(tmp_path) -> None:
    p = tmp_path / "a.txt"
    p.write_bytes("naïve résumé".encode("latin-1"))

    lw = LightweightExtractor()
    result = lw.extract(p, "text/plain")
    assert "naïve résumé" == result.text


def test_lightweight_extracts_md(tmp_path) -> None:
    p = tmp_path / "a.md"
    p.write_text("# Header\n\nbody text", encoding="utf-8")

    lw = LightweightExtractor()
    result = lw.extract(p, "text/markdown")
    assert "# Header" in result.text
    assert "body text" in result.text


def test_lightweight_extracts_html(tmp_path) -> None:
    p = tmp_path / "a.html"
    p.write_text(
        "<html><body><h1>title</h1><p>paragraph</p></body></html>",
        encoding="utf-8",
    )

    lw = LightweightExtractor()
    result = lw.extract(p, "text/html")
    # html2text emits markdown-ish; both strings should be present.
    assert "title" in result.text
    assert "paragraph" in result.text


def test_lightweight_extracts_csv(tmp_path) -> None:
    p = tmp_path / "a.csv"
    p.write_text("name,city\nalice,Berlin\nbob,Madrid\n", encoding="utf-8")

    lw = LightweightExtractor()
    result = lw.extract(p, "text/csv")
    assert "alice" in result.text
    assert "Berlin" in result.text


def test_lightweight_extracts_rtf(tmp_path) -> None:
    p = tmp_path / "a.rtf"
    p.write_text(
        r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Helvetica;}}"
        r"\f0\fs24 RTF body content here.\par}",
        encoding="ascii",
    )

    lw = LightweightExtractor()
    result = lw.extract(p, "application/rtf")
    assert "RTF body content here" in result.text


def test_lightweight_extracts_odt(tmp_path: Path) -> None:
    from odf.opendocument import OpenDocumentText
    from odf.text import P
    p = tmp_path / "a.odt"
    doc = OpenDocumentText()
    doc.text.addElement(P(text="odt paragraph one"))
    doc.text.addElement(P(text="odt paragraph two"))
    doc.save(str(p))

    lw = LightweightExtractor()
    result = lw.extract(p, "application/vnd.oasis.opendocument.text")
    assert "odt paragraph one" in result.text
    assert "odt paragraph two" in result.text


def test_lightweight_extracts_ics(tmp_path: Path) -> None:
    import datetime as dt
    from icalendar import Calendar, Event
    cal = Calendar()
    cal.add("prodid", "-//Test//Test//EN")
    cal.add("version", "2.0")
    ev = Event()
    ev.add("summary", "Annual review")
    ev.add("description", "Discuss quarterly bonus criteria")
    ev.add("location", "Conf room Berlin")
    ev.add("dtstart", dt.datetime(2026, 6, 1, 14, 0, tzinfo=dt.timezone.utc))
    cal.add_component(ev)

    p = tmp_path / "a.ics"
    p.write_bytes(cal.to_ical())

    lw = LightweightExtractor()
    result = lw.extract(p, "text/calendar")
    assert "Annual review" in result.text
    assert "quarterly bonus criteria" in result.text
    assert "Conf room Berlin" in result.text
    assert result.page_count == 1
    assert result.extractor == "lightweight@1.0"


def test_docling_import_warning_one_shot(caplog, monkeypatch) -> None:
    """When docling is missing, warn_docling_missing() emits exactly one
    WARN per process, even when called multiple times."""
    import localmail.search.extractor as ext_mod

    monkeypatch.setattr(ext_mod, "_DOCLING_WARNED", False, raising=False)
    monkeypatch.setattr(ext_mod, "_try_import_docling", lambda: None)

    with caplog.at_level(logging.WARNING, logger="localmail.search.extractor"):
        ext_mod.warn_docling_missing()
        ext_mod.warn_docling_missing()
        ext_mod.warn_docling_missing()

    warn_messages = [
        r for r in caplog.records
        if r.levelno == logging.WARNING
        and "extraction" in r.getMessage().lower()
    ]
    assert len(warn_messages) == 1


def test_docling_extractor_supports_pdf_only() -> None:
    """DoclingExtractor.supports() matches PDFs by MIME or extension only."""
    from localmail.search.extractor import DoclingExtractor
    de = DoclingExtractor()
    assert de.supports("application/pdf", "x.pdf")
    assert de.supports(None, "x.pdf")
    assert not de.supports("text/plain", "x.txt")
    assert not de.supports("image/png", "x.png")


def test_docling_extractor_raises_when_missing(monkeypatch, tmp_path) -> None:
    """If docling is not importable, .extract() raises ExtractorError."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import DoclingExtractor

    monkeypatch.setattr(ext_mod, "_try_import_docling", lambda: None)
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    de = DoclingExtractor()
    with pytest.raises(ExtractorError):
        de.extract(p, "application/pdf")


def test_docling_extractor_accepts_config() -> None:
    """DoclingExtractor.__init__ accepts a SearchConfig; defaults to
    SearchConfig() when omitted. The config is stored on the instance
    for the pipeline-options builder."""
    from localmail.config import SearchConfig
    from localmail.search.extractor import DoclingExtractor

    de = DoclingExtractor()
    assert de._cfg.extractor_docling_max_pages == 200
    assert de._cfg.extractor_ocr_languages == ["en"]

    custom = SearchConfig(
        extractor_docling_max_pages=50,
        extractor_ocr_languages=["en", "de", "ja"],
    )
    de2 = DoclingExtractor(custom)
    assert de2._cfg.extractor_docling_max_pages == 50
    assert de2._cfg.extractor_ocr_languages == ["en", "de", "ja"]


def test_docling_extractor_passes_configured_max_pages_to_convert(
    monkeypatch, tmp_path
) -> None:
    """The configured page cap reaches docling's convert() call.

    On docling 2.x the page limit is a convert-level argument
    (DocumentConverter.convert(..., max_num_pages=N)), not a
    PdfPipelineOptions field. This pins that extractor_docling_max_pages
    is actually forwarded so the OOM-guard cap takes effect.
    """
    import localmail.search.extractor as ext_mod
    from localmail.config import SearchConfig
    from localmail.search.extractor import DoclingExtractor

    recorded: dict[str, object] = {}

    class _FakeDoc:
        pages = [object()]

        def export_to_markdown(self) -> str:
            return "hello"

    class _FakeResult:
        document = _FakeDoc()

    class _FakeConverter:
        def __init__(self, *args: object, **kwargs: object) -> None:
            pass

        def convert(self, source: str, **kwargs: object) -> _FakeResult:
            recorded.update(kwargs)
            recorded["source"] = source
            return _FakeResult()

    monkeypatch.setattr(ext_mod, "_try_import_docling", lambda: _FakeConverter)

    pdf = tmp_path / "x.pdf"
    pdf.write_bytes(b"%PDF-1.4 dummy")

    de = DoclingExtractor(SearchConfig(extractor_docling_max_pages=50))
    de.extract(pdf, "application/pdf")

    assert recorded.get("max_num_pages") == 50


def test_iter_exc_chain_yields_self_then_cause() -> None:
    """The chain walk yields the exception, then its __cause__, in order."""
    from localmail.search.extractor import iter_exc_chain

    try:
        try:
            raise ValueError("root")
        except ValueError as inner:
            raise RuntimeError("wrap") from inner
    except RuntimeError as exc:
        chain = list(iter_exc_chain(exc))
    assert [type(e) for e in chain] == [RuntimeError, ValueError]


def test_iter_exc_chain_falls_back_to_context() -> None:
    """When __cause__ is None, the walk follows the implicit __context__."""
    from localmail.search.extractor import iter_exc_chain

    try:
        try:
            raise ValueError("root")
        except ValueError:
            raise RuntimeError("during handling")
    except RuntimeError as exc:
        chain = list(iter_exc_chain(exc))
    assert [type(e) for e in chain] == [RuntimeError, ValueError]


def test_iter_exc_chain_stops_on_suppress_context() -> None:
    """``raise X from None`` sets __suppress_context__ — the walk stops at X."""
    from localmail.search.extractor import iter_exc_chain

    try:
        try:
            raise ValueError("root")
        except ValueError:
            raise RuntimeError("masked") from None
    except RuntimeError as exc:
        chain = list(iter_exc_chain(exc))
    assert [type(e) for e in chain] == [RuntimeError]


# --- Third-party transient classification (#47) ------------------------------
#
# docling pulls models over the network (huggingface_hub) and raises
# third-party exception classes (requests / httpx / urllib3 / huggingface_hub)
# that are NOT in the builtin ConnectionError/TimeoutError hierarchy, so
# extract_worker._is_transient cannot recognise them. DoclingExtractor.extract
# opts these into TransientExtractorError so a model-download blip retries on
# the next sweep instead of poison-pilling a perfectly good PDF.


def _fake_converter_raising(exc_factory):
    """Build a fake docling ``DocumentConverter`` class whose ``convert()``
    raises ``exc_factory()``. Used to drive DoclingExtractor.extract without
    docling installed."""

    class _FakeConverter:
        def __init__(self, *args, **kwargs) -> None:
            pass

        def convert(self, source, **kwargs):
            raise exc_factory()

    return _FakeConverter


def _exc_class(name: str, module: str) -> type[Exception]:
    """Create a fresh Exception subclass whose ``__module__`` mimics a
    third-party package (e.g. ``requests.exceptions``)."""
    cls = type(name, (Exception,), {})
    cls.__module__ = module
    return cls


def test_exc_chain_has_transient_module_detects_top_level_package() -> None:
    """The pure helper matches on the top-level package, ignoring submodules."""
    from localmail.search.extractor import _exc_chain_has_transient_module

    httpx_err = _exc_class("ConnectError", "httpx._exceptions")
    assert _exc_chain_has_transient_module(httpx_err("refused"))


def test_exc_chain_has_transient_module_walks_cause_chain() -> None:
    """A transient third-party class anywhere in the cause chain matches."""
    from localmail.search.extractor import _exc_chain_has_transient_module

    hf_err = _exc_class("HfHubHTTPError", "huggingface_hub.errors")
    try:
        try:
            raise hf_err("502 fetching model")
        except Exception as inner:
            raise RuntimeError("docling pipeline failed") from inner
    except RuntimeError as exc:
        assert _exc_chain_has_transient_module(exc)


def test_exc_chain_has_transient_module_rejects_unknown_module() -> None:
    """A builtin / unknown-module exception is not opted in here (the builtin
    transient classes are extract_worker._is_transient's job)."""
    from localmail.search.extractor import _exc_chain_has_transient_module

    assert not _exc_chain_has_transient_module(ValueError("malformed bytes"))


def test_exc_chain_has_transient_module_respects_suppress_context() -> None:
    """``raise X from None`` stops the walk, matching _is_transient."""
    from localmail.search.extractor import _exc_chain_has_transient_module

    req_err = _exc_class("ConnectionError", "requests.exceptions")
    try:
        try:
            raise req_err("would otherwise be transient")
        except Exception:
            raise RuntimeError("deliberately masked") from None
    except RuntimeError as exc:
        assert not _exc_chain_has_transient_module(exc)


def test_docling_extract_classifies_requests_connection_error_transient(
    monkeypatch, tmp_path
) -> None:
    """A requests.exceptions.ConnectionError from docling.convert (model
    download blip) is raised as TransientExtractorError, not ExtractorError."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import (
        DoclingExtractor,
        TransientExtractorError,
    )

    req_err = _exc_class("ConnectionError", "requests.exceptions")
    monkeypatch.setattr(
        ext_mod,
        "_try_import_docling",
        lambda: _fake_converter_raising(lambda: req_err("max retries exceeded")),
    )
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    de = DoclingExtractor()
    with pytest.raises(TransientExtractorError):
        de.extract(p, "application/pdf")


def test_docling_extract_classifies_transient_in_cause_chain(
    monkeypatch, tmp_path
) -> None:
    """A generic docling exception whose __cause__ is a huggingface_hub error
    is still classified transient via the cause-chain walk."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import (
        DoclingExtractor,
        TransientExtractorError,
    )

    hf_err = _exc_class("HfHubHTTPError", "huggingface_hub.errors")

    def _raise_wrapped():
        try:
            raise hf_err("502 fetching model")
        except Exception as inner:
            raise RuntimeError("docling pipeline failed") from inner

    monkeypatch.setattr(
        ext_mod,
        "_try_import_docling",
        lambda: _fake_converter_raising(_raise_wrapped),
    )
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    de = DoclingExtractor()
    with pytest.raises(TransientExtractorError):
        de.extract(p, "application/pdf")


def test_docling_extract_keeps_value_error_as_permanent(
    monkeypatch, tmp_path
) -> None:
    """A genuine parse failure (ValueError, no transient module in the chain)
    stays a permanent ExtractorError — NOT TransientExtractorError — so a
    corrupt PDF is poison-pilled after retries."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import (
        DoclingExtractor,
        TransientExtractorError,
    )

    monkeypatch.setattr(
        ext_mod,
        "_try_import_docling",
        lambda: _fake_converter_raising(lambda: ValueError("corrupt PDF")),
    )
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    de = DoclingExtractor()
    with pytest.raises(ExtractorError) as excinfo:
        de.extract(p, "application/pdf")
    assert not isinstance(excinfo.value, TransientExtractorError)


# --- OCR engine selection + missing-engine classification (#248) --------------
#
# DoclingExtractor used to hardcode ocr_options=EasyOcrOptions(...). EasyOCR is
# not a docling dependency, so on an install without it every scanned PDF raised
# ImportError out of convert() on the *poison-pill* path, burning retry_count
# until the blob was given up on. 743 such rows on the live Mac archive. The
# hardcoding also overrode docling's own OcrAutoOptions default, which degrades
# to no-OCR without raising.


class _FakeOcrOptions:
    """Stand-in for a docling OCR options object."""

    def __init__(self, kind: str, lang: list[str]) -> None:
        self.kind = kind
        self.lang = lang


class _FakeOcrFactory:
    """Stand-in for docling's OCR factory, so engine-selection is testable
    without docling installed."""

    def __init__(self, kinds: list[str]) -> None:
        self.registered_kind = kinds
        self.created: list[tuple[str, list[str]]] = []

    def create_options(self, kind: str, **kwargs: object) -> _FakeOcrOptions:
        lang = kwargs.get("lang") or []
        self.created.append((kind, list(lang)))  # type: ignore[arg-type]
        return _FakeOcrOptions(kind, list(lang))  # type: ignore[arg-type]


def test_resolve_ocr_options_asks_the_factory_for_the_configured_kind(
    monkeypatch,
) -> None:
    """The config value is docling's own registry key — no mapping table."""
    import localmail.search.extractor as ext_mod
    from localmail.ocr_policy import plan_ocr

    factory = _FakeOcrFactory(["auto", "easyocr", "ocrmac"])
    monkeypatch.setattr(ext_mod, "_try_import_ocr_factory", lambda: factory)

    opts = ext_mod._resolve_ocr_options(plan_ocr("ocrmac"), ["en", "de"])

    assert factory.created == [("ocrmac", ["en", "de"])]
    assert opts is not None and opts.kind == "ocrmac"


def test_resolve_ocr_options_returns_none_when_ocr_is_disabled(
    monkeypatch,
) -> None:
    """engine='none' must not resolve an engine at all — do_ocr=False alone."""
    import localmail.search.extractor as ext_mod
    from localmail.ocr_policy import plan_ocr

    factory = _FakeOcrFactory(["auto"])
    monkeypatch.setattr(ext_mod, "_try_import_ocr_factory", lambda: factory)

    assert ext_mod._resolve_ocr_options(plan_ocr("none"), ["en"]) is None
    assert factory.created == []


def test_resolve_ocr_options_returns_none_when_the_factory_is_unavailable(
    monkeypatch,
) -> None:
    """An older docling without the factory must fall back to its own default
    rather than raise — the pre-#248 code had the same tolerance."""
    import localmail.search.extractor as ext_mod
    from localmail.ocr_policy import plan_ocr

    monkeypatch.setattr(ext_mod, "_try_import_ocr_factory", lambda: None)

    assert ext_mod._resolve_ocr_options(plan_ocr("easyocr"), ["en"]) is None


def test_an_unknown_ocr_engine_is_a_configuration_error(monkeypatch) -> None:
    """A typo names itself and lists the valid kinds — and must NOT be a
    poison pill, since no blob is at fault."""
    import localmail.search.extractor as ext_mod
    from localmail.ocr_policy import plan_ocr
    from localmail.search.extractor import (
        ExtractorConfigurationError,
        TransientExtractorError,
    )

    monkeypatch.setattr(
        ext_mod, "_try_import_ocr_factory", lambda: _FakeOcrFactory(["auto", "easyocr"])
    )

    with pytest.raises(ExtractorConfigurationError) as excinfo:
        ext_mod._resolve_ocr_options(plan_ocr("easyocrr"), ["en"])

    assert isinstance(excinfo.value, TransientExtractorError), (
        "a configuration error must never burn failed_extractions.retry_count"
    )
    msg = str(excinfo.value)
    assert "'easyocrr'" in msg and "auto, easyocr" in msg


def test_a_missing_ocr_engine_package_is_a_configuration_error(
    monkeypatch, tmp_path
) -> None:
    """The #248 failure itself: docling raises ImportError('EasyOCR is not
    installed...') from convert(). That is an install problem, not a bad blob,
    so it must classify as a configuration error and leave retry_count alone."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import (
        DoclingExtractor,
        ExtractorConfigurationError,
    )

    def _raise_missing_easyocr():
        raise ImportError(
            "EasyOCR is not installed. Please install it via `pip install easyocr` "
            "to use this OCR engine."
        )

    monkeypatch.setattr(
        ext_mod,
        "_try_import_docling",
        lambda: _fake_converter_raising(_raise_missing_easyocr),
    )
    p = tmp_path / "scan.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    with pytest.raises(ExtractorConfigurationError):
        DoclingExtractor().extract(p, "application/pdf")


def test_a_missing_ocr_engine_deep_in_the_cause_chain_still_classifies(
    monkeypatch, tmp_path
) -> None:
    """docling wraps engine construction, so the ImportError arrives as a
    __cause__ rather than at the top."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import (
        DoclingExtractor,
        ExtractorConfigurationError,
    )

    def _raise_wrapped():
        try:
            raise ImportError("RapidOCR is not installed.")
        except ImportError as inner:
            raise RuntimeError("pipeline build failed") from inner

    monkeypatch.setattr(
        ext_mod,
        "_try_import_docling",
        lambda: _fake_converter_raising(_raise_wrapped),
    )
    p = tmp_path / "scan.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    with pytest.raises(ExtractorConfigurationError):
        DoclingExtractor().extract(p, "application/pdf")


def test_a_corrupt_pdf_is_still_a_poison_pill_not_a_configuration_error(
    monkeypatch, tmp_path
) -> None:
    """The config-error branch must not swallow genuine blob failures — those
    still need to burn retry_count and land in failed_extractions."""
    import localmail.search.extractor as ext_mod
    from localmail.search.extractor import (
        DoclingExtractor,
        TransientExtractorError,
    )

    monkeypatch.setattr(
        ext_mod,
        "_try_import_docling",
        lambda: _fake_converter_raising(lambda: ValueError("corrupt PDF")),
    )
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF-1.4 dummy")

    with pytest.raises(ExtractorError) as excinfo:
        DoclingExtractor().extract(p, "application/pdf")
    assert not isinstance(excinfo.value, TransientExtractorError)


def test_warn_ocr_engine_unavailable_is_one_shot_per_process(caplog) -> None:
    """The daemon re-attempts every scanned PDF in the archive; one WARNING per
    process, not one per blob. Mirrors warn_docling_missing()."""
    import localmail.search.extractor as ext_mod

    ext_mod._OCR_ENGINE_WARNED = False
    with caplog.at_level(logging.WARNING, logger="localmail.search.extractor"):
        ext_mod.warn_ocr_engine_unavailable("easyocr", "EasyOCR is not installed.")
        ext_mod.warn_ocr_engine_unavailable("easyocr", "EasyOCR is not installed.")

    hits = [r for r in caplog.records if "easyocr" in r.getMessage()]
    assert len(hits) == 1
    assert "extractor_ocr_engine" in hits[0].getMessage()


def test_pdf_pipeline_options_reflect_the_configured_engine() -> None:
    """End-to-end against real docling: the built pipeline options carry the
    plan's do_ocr and the engine class the factory resolved. This is what
    #248 got wrong — do_ocr=True with a hardcoded EasyOcrOptions."""
    pytest.importorskip("docling")
    from localmail.config import SearchConfig
    from localmail.search.extractor import _build_pdf_pipeline_options

    auto = _build_pdf_pipeline_options(SearchConfig(extractor_ocr_engine="auto"))
    assert auto is not None
    assert auto.do_ocr is True
    assert type(auto.ocr_options).__name__ == "OcrAutoOptions"

    off = _build_pdf_pipeline_options(SearchConfig(extractor_ocr_engine="none"))
    assert off is not None
    assert off.do_ocr is False
