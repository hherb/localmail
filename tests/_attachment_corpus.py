"""Synthetic attachment fixture corpus for Phase 2 acceptance.

Builds in-memory bytes for each fixture format (no .eml or .pdf files
checked into the repo per CLAUDE.md). The acceptance harness writes
these bytes into a temp attachments_root tree and inserts
attachment_blobs + messages rows.

Each successful fixture is wrapped in a synthetic email whose subject
and body deliberately do NOT mention the attachment's distinctive
content. The attachment text contains a unique tag that is the target
of attachment_queries.json.

Negative fixtures (encrypted, corrupt, empty, oversized) have tag=None
so the harness knows to skip them in retrieval-quality scoring — they
are used only for the extraction-success gate.
"""

from __future__ import annotations

import datetime as _dt
import hashlib
import io
import json
from pathlib import Path
from typing import Any

_BASE = _dt.datetime(2026, 1, 1, tzinfo=_dt.timezone.utc)


# --- per-format builders ---------------------------------------------------


def build_native_pdf(text: str) -> bytes:
    """Return a native (text-layer) PDF containing *text* on the first page.

    Uses reportlab's canvas so the PDF has a real text stream that pypdf
    and pdfplumber can extract without OCR.
    """
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(72, 720, text)
    c.showPage()
    c.save()
    return buf.getvalue()


def build_scanned_pdf(text: str) -> bytes:
    """Return a PDF whose only content is a rasterized image of *text*.

    pypdf returns empty text on this; docling OCRs it. Used to verify
    the OCR fallback path in the extract worker.
    """
    from PIL import Image, ImageDraw
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    img = Image.new("RGB", (1000, 100), "white")
    ImageDraw.Draw(img).text((10, 30), text, fill="black")
    img_buf = io.BytesIO()
    img.save(img_buf, format="PNG")
    img_buf.seek(0)

    pdf_buf = io.BytesIO()
    c = canvas.Canvas(pdf_buf)
    c.drawImage(ImageReader(img_buf), 72, 600, width=400, height=80)
    c.showPage()
    c.save()
    return pdf_buf.getvalue()


def build_docx(paragraphs: list[str]) -> bytes:
    """Return a .docx blob containing *paragraphs* as separate paragraph runs."""
    import docx

    buf = io.BytesIO()
    d = docx.Document()
    for p in paragraphs:
        d.add_paragraph(p)
    d.save(buf)
    return buf.getvalue()


def build_xlsx(sheets: dict[str, list[list[str]]]) -> bytes:
    """Return a .xlsx blob where each key in *sheets* maps to a worksheet of rows.

    The openpyxl default 'Sheet' is removed when it is not among the
    requested sheet names.
    """
    import openpyxl

    wb = openpyxl.Workbook()
    default_name = wb.sheetnames[0]
    for sheet_name, rows in sheets.items():
        if sheet_name == default_name:
            ws = wb[default_name]
        else:
            ws = wb.create_sheet(title=sheet_name)
        for row in rows:
            ws.append(row)
    if default_name not in sheets and default_name in wb.sheetnames:
        del wb[default_name]
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_pptx(title: str, notes: str) -> bytes:
    """Return a .pptx blob with one blank slide whose speaker notes contain *notes*."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.notes_slide.notes_text_frame.text = notes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_rtf(text: str) -> bytes:
    """Return a minimal RTF document containing *text* in Helvetica 12pt."""
    body = (
        r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Helvetica;}}"
        r"\f0\fs24 " + text + r"\par}"
    )
    return body.encode("ascii", errors="ignore")


def build_txt(text: str, encoding: str = "utf-8") -> bytes:
    """Return plain text bytes encoded with *encoding* (default UTF-8)."""
    return text.encode(encoding)


def build_html(title: str, body: str) -> bytes:
    """Return a minimal HTML document with *title* in an h1 and *body* in a p."""
    return (
        f"<html><body><h1>{title}</h1><p>{body}</p></body></html>"
    ).encode("utf-8")


def build_csv(rows: list[list[str]]) -> bytes:
    """Return CSV bytes for the given *rows* (list of lists)."""
    import csv

    buf = io.StringIO()
    w = csv.writer(buf)
    for r in rows:
        w.writerow(r)
    return buf.getvalue().encode("utf-8")


def build_ics(summary: str, description: str, location: str) -> bytes:
    """Return iCalendar bytes for a single event with the given fields."""
    from icalendar import Calendar, Event

    cal = Calendar()
    cal.add("prodid", "-//Test//Test//EN")
    cal.add("version", "2.0")
    ev = Event()
    ev.add("summary", summary)
    ev.add("description", description)
    ev.add("location", location)
    ev.add("dtstart", _BASE)
    cal.add_component(ev)
    return cal.to_ical()


def build_encrypted_pdf() -> bytes:
    """Return a password-protected PDF (owner='o', user='u', R=4 encryption).

    LightweightExtractor should raise ExtractorError on this; the extract
    worker lands it in failed_extractions. Useful as a negative-test fixture.
    """
    import pikepdf

    src = build_native_pdf("secret content")
    out = io.BytesIO()
    with pikepdf.open(io.BytesIO(src)) as p:
        p.save(out, encryption=pikepdf.Encryption(owner="o", user="u", R=4))
    return out.getvalue()


def build_corrupt_pdf() -> bytes:
    """Return a byte sequence that begins with a PDF magic number but is otherwise invalid.

    Parsers that attempt to read the cross-reference table or content
    streams will fail. Used as a negative-test fixture.
    """
    return b"%PDF-1.4\nthis is not a valid PDF body"


def build_empty() -> bytes:
    """Return zero bytes — the smallest possible 'attachment'.

    The extractor should treat this as an error or empty extraction.
    Used as a negative-test fixture.
    """
    return b""


def build_oversized_pdf(min_bytes: int) -> bytes:
    """Return a PDF padded with comment lines until its size exceeds *min_bytes*.

    Used to verify that the extract worker respects its configured size cap
    (default 50 MB). The padding consists of `%pad` comment lines which are
    valid PDF syntax and do not break parsers.
    """
    base = build_native_pdf("a")
    padding = b"%pad\n" * ((min_bytes // 5) + 1)
    return base + padding


# --- in-DB seeding ---------------------------------------------------------


def _write_blob(content: bytes, attachments_root: Path) -> tuple[bytes, Path]:
    """Write *content* to the content-addressable blob tree under *attachments_root*.

    Returns the raw SHA-256 digest (32 bytes) and the absolute blob path.
    If the blob already exists it is not overwritten.
    """
    sha = hashlib.sha256(content).digest()
    sub = sha.hex()
    blob_path = attachments_root / "blobs" / sub[:2] / sub[2:4] / sub
    blob_path.parent.mkdir(parents=True, exist_ok=True)
    if not blob_path.exists():
        blob_path.write_bytes(content)
    return sha, blob_path


def build_corpus(
    conn: Any,
    *,
    attachments_root: Path,
) -> list[dict[str, Any]]:
    """Seed attachment fixtures into the test DB and write blobs to disk.

    Creates one account named 'attach_corpus', then inserts:
      - 11 fixture messages (one per allowlisted format), each with a
        unique *tag* string embedded in the attachment that does NOT
        appear in the wrapping email's subject or body.
      - 4 negative-test fixtures (encrypted PDF, corrupt PDF, empty
        blob, oversized PDF) with tag=None.
      - 5 noise messages with no attachments.

    Each fixture message gets one attachment_blobs row (keyed on sha256)
    and one messages row (with attachments JSONB referencing the blob).

    Returns a list of dicts with keys:
      id         — integer primary key from messages
      subject    — email subject
      tag        — distinctive string in the attachment (None for negatives)
      mime       — MIME type string
      filename   — original filename as stored in attachments JSONB

    Commits the transaction before returning.
    """
    fixtures: list[dict[str, Any]] = [
        {
            "subject": "Quarterly review attachment",
            "body": "see attached",
            "tag": "non-disclosure obligations under section 5",
            "filename": "contract.pdf",
            "mime": "application/pdf",
            "bytes": build_native_pdf("non-disclosure obligations under section 5"),
        },
        {
            "subject": "Scanned report",
            "body": "FYI scan",
            "tag": "annual revenue growth fourteen percent",
            "filename": "scan.pdf",
            "mime": "application/pdf",
            "bytes": build_scanned_pdf("annual revenue growth fourteen percent"),
        },
        {
            "subject": "Onboarding pack",
            "body": "details inside",
            "tag": "probation period three months from start date",
            "filename": "onboarding.docx",
            "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "bytes": build_docx([
                "Welcome to the team.",
                "probation period three months from start date",
            ]),
        },
        {
            "subject": "Budget spreadsheet",
            "body": "numbers attached",
            "tag": "marketing line item Berlin office Q4",
            "filename": "budget.xlsx",
            "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "bytes": build_xlsx({
                "Q4": [
                    ["Office", "Line item", "Amount"],
                    ["Berlin", "marketing line item Berlin office Q4", "12000"],
                ]
            }),
        },
        {
            "subject": "All-hands deck",
            "body": "slides",
            "tag": "company-wide hackathon scheduled mid October",
            "filename": "allhands.pptx",
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "bytes": build_pptx(
                "All-Hands",
                "company-wide hackathon scheduled mid October",
            ),
        },
        {
            "subject": "Notes",
            "body": "see notes",
            "tag": "research direction transformer attention sparsity",
            "filename": "notes.rtf",
            "mime": "application/rtf",
            "bytes": build_rtf("research direction transformer attention sparsity"),
        },
        {
            "subject": "Quick note",
            "body": "attached",
            "tag": "vendor SLA monthly recurring fee adjustment",
            "filename": "note.txt",
            "mime": "text/plain",
            "bytes": build_txt("vendor SLA monthly recurring fee adjustment"),
        },
        {
            "subject": "Doc",
            "body": "see md",
            "tag": "rollout plan Tuesday wave one customers",
            "filename": "doc.md",
            "mime": "text/markdown",
            "bytes": build_txt("# Rollout\n\nrollout plan Tuesday wave one customers"),
        },
        {
            "subject": "Web archive",
            "body": "attached html",
            "tag": "policy update privacy controls fourth quarter",
            "filename": "archive.html",
            "mime": "text/html",
            "bytes": build_html(
                "Policy",
                "policy update privacy controls fourth quarter",
            ),
        },
        {
            "subject": "Tabular data",
            "body": "csv attached",
            "tag": "compliance audit findings critical severity",
            "filename": "data.csv",
            "mime": "text/csv",
            "bytes": build_csv([
                ["finding_id", "severity", "summary"],
                ["F-1", "critical", "compliance audit findings critical severity"],
            ]),
        },
        {
            "subject": "Invite",
            "body": "calendar attached",
            "tag": "interview panel candidate ML researcher",
            "filename": "invite.ics",
            "mime": "text/calendar",
            "bytes": build_ics(
                "interview panel candidate ML researcher",
                "ML researcher screening",
                "Conf room A",
            ),
        },
    ]

    negatives: list[dict[str, Any]] = [
        {
            "subject": "Locked file",
            "body": "password protected",
            "tag": None,
            "filename": "locked.pdf",
            "mime": "application/pdf",
            "bytes": build_encrypted_pdf(),
        },
        {
            "subject": "Mystery attachment",
            "body": "broken",
            "tag": None,
            "filename": "broken.pdf",
            "mime": "application/pdf",
            "bytes": build_corrupt_pdf(),
        },
        {
            "subject": "Empty placeholder",
            "body": "empty",
            "tag": None,
            "filename": "empty.txt",
            "mime": "text/plain",
            "bytes": build_empty(),
        },
        {
            "subject": "Huge PDF",
            "body": "too big to extract",
            "tag": None,
            "filename": "huge.pdf",
            "mime": "application/pdf",
            "bytes": build_oversized_pdf(60 * 1024 * 1024),  # >50 MB default cap
        },
    ]

    noise: list[dict[str, str]] = [
        {"subject": "Project status sync", "body": "Recurring weekly sync agenda"},
        {"subject": "Reminder about timezone update", "body": "switch to CET"},
        {"subject": "Sourdough starter discard", "body": "kitchen tips"},
        {"subject": "Library closing hours", "body": "summer schedule"},
        {"subject": "Garden compost run", "body": "weekend"},
    ]

    seeded: list[dict[str, Any]] = []
    with conn.cursor() as cur:
        cur.execute(
            "INSERT INTO accounts (name, email_address, imap_host, auth_method) "
            "VALUES ('attach_corpus', 'a@b', 'h', 'password') RETURNING id"
        )
        row = cur.fetchone()
        assert row is not None
        acct = row[0]

        for i, f in enumerate(fixtures + negatives):
            sha, blob_path = _write_blob(f["bytes"], attachments_root)
            cur.execute(
                "INSERT INTO attachment_blobs (sha256, path, mime_type, size_bytes) "
                "VALUES (%s, %s, %s, %s) ON CONFLICT DO NOTHING",
                (sha, str(blob_path), f["mime"], len(f["bytes"])),
            )
            attachments_json = json.dumps([{
                "filename": f["filename"], "sha256": sha.hex()
            }])
            cur.execute(
                "INSERT INTO messages"
                " (account_id, message_id, raw_sha256, subject, body_text,"
                "  headers, raw_bytes, size_bytes, attachments, date_sent)"
                " VALUES (%s, %s, %s, %s, %s, '{}'::jsonb, %s, %s, %s::jsonb, %s)"
                " RETURNING id",
                (
                    acct,
                    f"<ac{i}@local>",
                    bytes([i + 1]) * 32,
                    f["subject"],
                    f["body"],
                    b"raw",
                    len(f["body"]),
                    attachments_json,
                    _BASE + _dt.timedelta(days=i),
                ),
            )
            mid_row = cur.fetchone()
            assert mid_row is not None
            seeded.append({
                "id": mid_row[0],
                "subject": f["subject"],
                "tag": f["tag"],
                "mime": f["mime"],
                "filename": f["filename"],
            })

        for j, n in enumerate(noise):
            cur.execute(
                "INSERT INTO messages"
                " (account_id, message_id, raw_sha256, subject, body_text,"
                "  headers, raw_bytes, size_bytes, attachments, date_sent)"
                " VALUES (%s, %s, %s, %s, %s, '{}'::jsonb, %s, %s, '[]'::jsonb, %s)",
                (
                    acct,
                    f"<noise{j}@local>",
                    bytes([j + 200]) * 32,
                    n["subject"],
                    n["body"],
                    b"raw",
                    len(n["body"]),
                    _BASE + _dt.timedelta(days=100 + j),
                ),
            )

    conn.commit()
    return seeded
