# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Horst Herb

"""Attachment extractors.

Protocol + LightweightExtractor (pure-Python, no OCR, covers 11 allowlisted
formats) + DoclingExtractor (lazy-imported, OCR-capable, added later).
The extract_worker picks LightweightExtractor by default; if it returns
empty/raises on a PDF, the worker falls back to DoclingExtractor when
docling is importable.

Exports:
- ExtractedText: frozen dataclass returned by every successful extraction.
- ExtractorError: raised on irrecoverable failure.
- AttachmentExtractor: Protocol that all extractors implement.
- LightweightExtractor: pure-Python dispatch across PDF/Office/text/ODT/ICS.
"""

from __future__ import annotations

import logging
from collections.abc import Iterator
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Protocol, cast, runtime_checkable

from localmail.config import SearchConfig
from localmail.ocr_policy import (
    OCR_DISABLED,
    OcrPlan,
    plan_ocr,
    unknown_engine_message,
)
from localmail.pgtext import strip_nuls
from localmail.search.attachment_kind import extension_of, is_pdf

_LOG = logging.getLogger(__name__)


@dataclass(frozen=True)
class ExtractedText:
    """The result of extracting text from a blob.

    ``text`` is normalised on construction to contain no NUL bytes — see
    ``__post_init__``.

    Attributes:
        text: The extracted plain-text, NUL-free. May be '' (sentinel meaning
            "we tried, got nothing, don't retry").
        page_count: Optional logical page count (PDFs, Office docs). None
            for plain-text formats like TXT/MD/HTML/CSV/ICS.
        extractor: Identifier of the extractor that produced the text,
            including version. Values used elsewhere in the codebase:
            'lightweight@1.0', 'docling@X.Y', 'lightweight-empty',
            'size-skipped'.
    """

    text: str
    page_count: int | None
    extractor: str

    def __post_init__(self) -> None:
        """Strip NUL bytes from ``text`` so the value is safe to INSERT (#249).

        Postgres TEXT rejects ``\\x00`` and ``attachment_text.extracted_text``
        is this type's only consumer, so a NUL that survives to the INSERT
        aborts it. That abort escaped ``_process_blob`` into the worker's outer
        safety net, which recorded the blob in ``failed_extractions`` under the
        extractor name ``'unexpected'`` — and because the same bytes always
        re-extract to the same NUL, the retry budget was spent on a failure no
        retry could clear.

        Extracted text inherits whatever the source document contained, so this
        is not exotic: 128 blobs on the live Mac archive (112 PDFs, 10
        text/plain, 5 octet-stream, 1 html) had been given up on this way.

        Normalising here rather than in each of the eleven ``_extract_*``
        methods means a twelfth cannot forget — the same by-construction
        reasoning as ``open_attachment_bytes``' unconditional ACL check (#67).
        ``object.__setattr__`` is the standard idiom for a frozen dataclass;
        the class stays frozen to every other caller. ``strip_nuls`` returns
        the original object on the clean path, so the common case rebinds the
        same string rather than copying it.
        """
        object.__setattr__(self, "text", strip_nuls(self.text))


class ExtractorError(Exception):
    """Raised by an extractor on irrecoverable failure.

    The caller (extract_worker) records the blob in failed_extractions
    and continues processing the next blob in the batch.
    """


class TransientExtractorError(ExtractorError):
    """Raised by an extractor when the failure is *not* the blob's fault.

    Transient conditions — model-download blips, OCR backend OOM, retried-
    once IO hiccups — must not increment a blob's ``failed_extractions``
    retry_count. The worker classifies these via ``_is_transient`` and
    rolls back the SAVEPOINT without recording, leaving the blob eligible
    for the next sweep.

    Extractors may raise this directly when they detect the cause is
    transient. Otherwise the worker falls back to walking the exception's
    cause chain for built-in transient classes (ConnectionError,
    TimeoutError, MemoryError).
    """


class ExtractorConfigurationError(TransientExtractorError):
    """Raised when the extractor cannot run because of how it is *configured*,
    not because of anything wrong with the blob (#248).

    Two causes today, both about OCR: an engine name docling does not register,
    and an engine whose package is not installed (docling raises ``ImportError``
    out of ``convert()``).

    Subclassing ``TransientExtractorError`` is the load-bearing part: the worker
    already treats that as "not the blob's fault", so ``retry_count`` is never
    burned and no ``failed_extractions`` row is written. Burning it *was* #248 —
    EasyOCR-not-installed landed on the poison-pill path, so every scanned PDF in
    the archive accumulated three failures and was then given up on, for a
    problem no retry could fix and that had nothing to do with the document.
    Scanned PDFs are precisely what the docling fallback exists for.

    The bound is therefore the *transient* budget
    (``extract_worker_max_transient_retries``, #153), which exists for exactly
    this shape: a not-the-blob's-fault failure that may nonetheless be permanent.
    A dedicated ``attachment_text`` sentinel was rejected — it would make the
    blob ineligible for re-claim, so fixing the config would silently *not*
    re-open the documents it was fixed for (the one-way door ``type-skipped``
    documents). Recovery is ``localmail retry-failed-extractions``, which clears
    both tables.
    """


@runtime_checkable
class AttachmentExtractor(Protocol):
    """The contract every attachment extractor must implement.

    `@runtime_checkable` lets the extract_worker verify a configured
    extractor instance satisfies the protocol via isinstance() at startup.
    """

    name: str
    version: str

    def supports(self, mime_type: str | None, filename: str | None) -> bool:
        """Return True iff this extractor can process the given blob.

        Implementations should accept either a known MIME type OR a
        matching filename extension (mail clients commonly mis-set MIME).
        """
        ...

    def extract(
        self, blob_path: Path, mime_type: str | None, *, filename: str | None = None
    ) -> ExtractedText:
        """Extract text from the blob at `blob_path`.

        `filename` is the blob's *original* per-message name, used for format
        dispatch when the MIME type is unhelpful. It is a separate argument
        because `blob_path` is content-addressable
        (`blobs/<aa>/<bb>/<sha256hex>`) and carries no extension — reading one
        off it was #216.

        Returns an ExtractedText. May return ExtractedText(text='', ...)
        when the file is structurally valid but contains no extractable
        text (e.g. a scanned PDF for a non-OCR extractor). Raises
        ExtractorError on irrecoverable failure (parse error, encryption,
        truncated bytes, etc.).
        """
        ...


# --- Lightweight extractor ---------------------------------------------------
#
# Pure-Python, no OCR. The MIME/extension allowlists below are derived from
# SearchConfig defaults so there is a single source of truth.

# Derived from SearchConfig defaults — single source of truth.
_LW_MIME_PREFIXES: frozenset[str] = frozenset(
    cast(Callable[[], list[str]],
         SearchConfig.model_fields["extractor_mime_allowlist"].default_factory)()
)
_LW_EXTENSIONS: frozenset[str] = frozenset(
    cast(Callable[[], list[str]],
         SearchConfig.model_fields["extractor_extension_allowlist"].default_factory)()
)


class LightweightExtractor:
    """Pure-Python extractor for documents that don't require OCR.

    Dispatches across all 11 allowlisted formats via extract():
    PDF (pypdf), DOCX (python-docx), XLSX (openpyxl), PPTX (python-pptx),
    TXT/MD (chardet fallback), HTML (html2text), CSV (stdlib csv),
    RTF (striprtf), ODT (odfpy), ICS (icalendar).
    """

    name = "lightweight"
    version = "1.0"

    def __init__(self, cfg: "SearchConfig | None" = None) -> None:
        """Construct with an optional SearchConfig.

        If cfg is None, defaults to SearchConfig(). Configs are only
        read for numeric tunables (e.g. chardet confidence threshold) —
        extractor behaviour otherwise stays config-agnostic.
        """
        self._cfg = cfg if cfg is not None else SearchConfig()

    def supports(self, mime_type: str | None, filename: str | None) -> bool:
        """True iff the MIME type or filename extension is allowlisted."""
        if mime_type and mime_type.lower() in _LW_MIME_PREFIXES:
            return True
        return extension_of(filename) in _LW_EXTENSIONS

    def extract(
        self, blob_path: Path, mime_type: str | None, *, filename: str | None = None
    ) -> ExtractedText:
        """Extract text from `blob_path`, dispatching on `mime_type` or on
        `filename`'s extension.

        The extension comes from `filename`, never from `blob_path`: the blob
        path is content-addressable and extensionless, so dispatching on it left
        every mis-typed attachment falling through to `ExtractorError` (#216).

        Returns ExtractedText with text='' when the file is structurally valid
        but contains no extractable text (e.g. a scanned PDF with no native text
        stream). Raises ExtractorError on irrecoverable parse failures (corrupt
        bytes, encryption, etc.).
        """
        ext = extension_of(filename)
        mt = (mime_type or "").lower()

        if mt == "application/pdf" or ext == ".pdf":
            return self._extract_pdf(blob_path)
        if "wordprocessingml" in mt or ext == ".docx":
            return self._extract_docx(blob_path)
        if "spreadsheetml" in mt or ext == ".xlsx":
            return self._extract_xlsx(blob_path)
        if "presentationml" in mt or ext == ".pptx":
            return self._extract_pptx(blob_path)
        if mt == "text/plain" or ext == ".txt":
            return self._extract_txt(blob_path)
        if mt == "text/markdown" or ext == ".md":
            return self._extract_md(blob_path)
        if mt == "text/html" or ext in (".html", ".htm"):
            return self._extract_html(blob_path)
        if mt == "text/csv" or ext == ".csv":
            return self._extract_csv(blob_path)
        if mt == "application/rtf" or ext == ".rtf":
            return self._extract_rtf(blob_path)
        if "opendocument.text" in mt or ext == ".odt":
            return self._extract_odt(blob_path)
        if mt == "text/calendar" or ext == ".ics":
            return self._extract_ics(blob_path)

        raise ExtractorError(f"no extractor for {mt!r}/{ext!r}")

    def _extract_pdf(self, blob_path: Path) -> ExtractedText:
        """Extract text from a PDF blob using pypdf.

        Returns ExtractedText(text='', ...) on PDFs whose only content
        is rasterized images (scanned documents). Raises ExtractorError
        on encrypted/password-protected PDFs and on irrecoverable parse
        failures.
        """
        # Logging deferred to extract_worker (Task 13): it catches
        # ExtractorError, records the blob in failed_extractions, and
        # emits the log line. Keeping this method silent matches the
        # Phase 1 chunking.py / embed_worker.py separation of concerns.
        import pypdf
        try:
            reader = pypdf.PdfReader(str(blob_path))
        except Exception as exc:
            raise ExtractorError(f"pypdf failed to open: {exc}") from exc

        if reader.is_encrypted:
            raise ExtractorError(
                "pypdf: encrypted PDF (no password supplied)"
            )

        try:
            pages = [page.extract_text() or "" for page in reader.pages]
        except Exception as exc:
            raise ExtractorError(f"pypdf failed to extract: {exc}") from exc

        text = "\n".join(pages).strip()
        return ExtractedText(
            text=text,
            page_count=len(reader.pages),
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_docx(self, blob_path: Path) -> ExtractedText:
        """Extract text from a DOCX blob using python-docx.

        Logging deferred to extract_worker (Task 13): silent on success
        and on raise; the worker catches ExtractorError and reports.
        """
        import docx
        try:
            d = docx.Document(str(blob_path))
        except Exception as exc:
            raise ExtractorError(f"python-docx failed to open: {exc}") from exc
        paras = [p.text for p in d.paragraphs if p.text]
        text = "\n".join(paras).strip()
        return ExtractedText(
            text=text, page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_xlsx(self, blob_path: Path) -> ExtractedText:
        """Extract text from an XLSX blob using openpyxl.

        Reads all sheets; each row is joined into a single text line.
        Logging deferred to extract_worker (Task 13).

        Note: blob paths have no extension (SHA-256 hex names), so we pass
        the bytes via io.BytesIO rather than the file path — openpyxl uses
        the filename extension to detect the format when given a path string,
        which would raise an error for extensionless blob files.
        """
        import contextlib
        import io
        import openpyxl
        try:
            wb = openpyxl.load_workbook(
                io.BytesIO(blob_path.read_bytes()), read_only=True, data_only=True
            )
        except Exception as exc:
            raise ExtractorError(f"openpyxl failed to open: {exc}") from exc

        parts: list[str] = []
        with contextlib.closing(wb):
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text:
                        parts.append(row_text)
        text = "\n".join(parts).strip()
        return ExtractedText(
            text=text, page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_pptx(self, blob_path: Path) -> ExtractedText:
        """Extract text from a PPTX blob using python-pptx.

        Captures shape-frame text (including field elements like slide
        numbers and date fields) plus speaker notes. Logging deferred
        to extract_worker (Task 13).
        """
        from pptx import Presentation
        try:
            prs = Presentation(str(blob_path))
        except Exception as exc:
            raise ExtractorError(f"python-pptx failed to open: {exc}") from exc

        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        parts.append(t)
            if slide.has_notes_slide:
                ntf = slide.notes_slide.notes_text_frame
                if ntf is not None:
                    notes = ntf.text
                    if notes:
                        parts.append(notes)

        text = "\n".join(parts).strip()
        return ExtractedText(
            text=text, page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_txt(self, blob_path: Path) -> ExtractedText:
        """Extract text from a plain-text blob.

        Encoding detection order:
          1. UTF-8 (strict) — covers the vast majority of modern files.
          2. chardet heuristic — handles multi-byte encodings (Shift-JIS,
             GB18030, etc.) that are unambiguously detectable by chardet.
             A CHARDET_CONFIDENCE_MIN threshold guards against low-confidence
             guesses where chardet cannot distinguish between similar
             single-byte encodings (e.g. latin-1 vs cp1250 for short files).
          3. latin-1 hard fallback — IANA-registered default encoding for
             text; maps all 256 byte values without error, so this path
             never raises UnicodeDecodeError.

        errors='replace' on the chardet path absorbs any remaining
        mis-detections without raising. Logging deferred to extract_worker
        (Task 13).
        """
        import chardet

        raw = blob_path.read_bytes()
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            det = chardet.detect(raw) or {}
            conf = det.get("confidence") or 0.0
            enc = (
                det.get("encoding")
                if conf >= self._cfg.extractor_chardet_confidence_min
                else None
            )
            encoding = enc or "latin-1"
            try:
                text = raw.decode(encoding, errors="replace")
            except Exception as exc:
                raise ExtractorError(f"txt decode failed: {exc}") from exc
        return ExtractedText(
            text=text.strip(),
            page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_md(self, blob_path: Path) -> ExtractedText:
        """Extract text from a Markdown blob.

        Markdown is left as-is — chunking and embeddings handle its
        structure well without stripping. Delegates to _extract_txt for
        encoding detection (DRY: Markdown is plain text with structure
        markers we don't interpret). Logging deferred to extract_worker.
        """
        return self._extract_txt(blob_path)

    def _extract_html(self, blob_path: Path) -> ExtractedText:
        """Extract text from an HTML blob using html2text.

        Returns Markdown-ish output: markup is stripped, paragraphs and
        headings are preserved as plain text. ignore_images drops
        '![alt](url)' noise; body_width=0 disables line-wrapping so full
        sentences stay on one line for better snippet quality. Logging
        deferred to extract_worker (Task 13).
        """
        import html2text
        try:
            html = blob_path.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            raise ExtractorError(f"html read failed: {exc}") from exc
        h = html2text.HTML2Text()
        h.ignore_images = True
        h.body_width = 0
        text = h.handle(html).strip()
        return ExtractedText(
            text=text, page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_csv(self, blob_path: Path) -> ExtractedText:
        """Extract text from a CSV blob using stdlib csv.

        Each row's cells are space-joined; rows are newline-joined. The
        newline='' argument lets csv.reader handle cross-platform line
        endings inside quoted cells correctly. Logging deferred to
        extract_worker (Task 13).
        """
        import csv
        try:
            with blob_path.open(
                "r", encoding="utf-8", errors="replace", newline=""
            ) as f:
                rows = [" ".join(r) for r in csv.reader(f)]
        except Exception as exc:
            raise ExtractorError(f"csv read failed: {exc}") from exc
        return ExtractedText(
            text="\n".join(rows).strip(),
            page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_rtf(self, blob_path: Path) -> ExtractedText:
        """Extract text from an RTF blob using striprtf.

        Reads as UTF-8 with replacement (RTF files are nominally ASCII
        with backslash-escaped non-ASCII; replacement chars are harmless
        for the rare malformed case). Logging deferred to extract_worker
        (Task 13).
        """
        from striprtf.striprtf import rtf_to_text
        try:
            raw = blob_path.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            raise ExtractorError(f"rtf read failed: {exc}") from exc
        try:
            text = rtf_to_text(raw)
        except Exception as exc:
            raise ExtractorError(f"striprtf failed: {exc}") from exc
        return ExtractedText(
            text=text.strip(),
            page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_odt(self, blob_path: Path) -> ExtractedText:
        """Extract text from an ODT blob using odfpy.

        Walks all <text:p> paragraph elements and concatenates their
        character data. Logging deferred to extract_worker (Task 13).
        """
        from odf.opendocument import load
        from odf.text import P
        try:
            doc = load(str(blob_path))
        except Exception as exc:
            raise ExtractorError(f"odfpy failed to load: {exc}") from exc
        paras = doc.getElementsByType(P)
        text = "\n".join(
            "".join(node.data for node in p.childNodes if hasattr(node, "data"))
            for p in paras
        ).strip()
        return ExtractedText(
            text=text, page_count=None,
            extractor=f"{self.name}@{self.version}",
        )

    def _extract_ics(self, blob_path: Path) -> ExtractedText:
        """Extract text from an ICS calendar blob using icalendar.

        For each VEVENT, concatenates SUMMARY, DESCRIPTION, LOCATION,
        DTSTART, and ATTENDEE values as text. page_count carries the
        event count (or None when no events). Logging deferred to
        extract_worker (Task 13).
        """
        from icalendar import Calendar
        try:
            raw = blob_path.read_bytes()
            cal = Calendar.from_ical(raw)
        except Exception as exc:
            raise ExtractorError(f"icalendar parse failed: {exc}") from exc

        parts: list[str] = []
        event_count = 0
        for component in cal.walk():
            if component.name == "VEVENT":
                event_count += 1
                for field in ("SUMMARY", "DESCRIPTION", "LOCATION"):
                    val = component.get(field)
                    if val:
                        parts.append(str(val))
                dtstart = component.get("DTSTART")
                if dtstart:
                    parts.append(str(dtstart.dt))
                attendees = component.get("ATTENDEE", [])
                if attendees:
                    # ATTENDEE may be a single value or a list
                    if not isinstance(attendees, list):
                        attendees = [attendees]
                    for attendee in attendees:
                        parts.append(str(attendee))

        return ExtractedText(
            text="\n".join(parts).strip(),
            page_count=event_count or None,
            extractor=f"{self.name}@{self.version}",
        )


# --- Docling (optional, OCR-capable) -----------------------------------------
#
# Docling is in the [extraction] uv extra — installed via
# `uv sync --extra extraction`. The extractor lazy-imports it on first
# call. When not installed, `warn_docling_missing()` emits exactly one
# WARN per process pointing at the install hint, and .extract() raises
# ExtractorError so the extract_worker records the failure cleanly.

_DOCLING_WARNED = False


def _try_import_docling() -> "type | None":
    """Return docling's `DocumentConverter` class, or `None` if docling
    is not installed.

    Indirected for test monkeypatching: tests replace this function to
    simulate the "missing" state without uninstalling docling.
    """
    try:
        from docling.document_converter import DocumentConverter
        return DocumentConverter
    except ImportError:
        return None


def warn_docling_missing() -> None:
    """Emit a one-shot WARN per process pointing at the docling install hint.

    The extract_worker calls this on the first PDF where
    lightweight extraction returned empty/raised AND docling is not
    importable. Subsequent calls in the same process are silent so a
    large archive sync doesn't flood the log.
    """
    global _DOCLING_WARNED
    if _DOCLING_WARNED:
        return
    _DOCLING_WARNED = True
    _LOG.warning(
        "docling is not installed; PDFs that lightweight cannot read "
        "will be marked as lightweight-empty. Install with "
        "`uv sync --extra extraction` to enable OCR fallback for "
        "scanned PDFs."
    )


_OCR_ENGINE_WARNED = False


def warn_ocr_engine_unavailable(engine: str, detail: str) -> None:
    """Emit a one-shot WARN per process naming the unusable OCR engine (#248).

    One line per *process*, not per blob: the daemon re-attempts every scanned
    PDF in the archive, and the pre-#248 behaviour flooded the log with the same
    EasyOCR message while quietly poison-pilling each document. Mirrors
    ``warn_docling_missing()``.
    """
    global _OCR_ENGINE_WARNED
    if _OCR_ENGINE_WARNED:
        return
    _OCR_ENGINE_WARNED = True
    _LOG.warning(
        "OCR engine %r (search.extractor_ocr_engine) is not usable: %s "
        "Scanned PDFs will not be indexed until this is fixed; their extraction "
        "is being held, not failed, so `localmail retry-failed-extractions` "
        "will pick them up afterwards. Set it to %r to accept that and stop "
        "trying, or to 'auto' to use whichever engine is installed.",
        engine,
        detail,
        OCR_DISABLED,
    )


def _try_import_ocr_factory() -> Any:
    """Return docling's OCR options factory, or ``None`` when unavailable.

    Typed ``Any`` rather than a docling class: docling lives in the optional
    ``[extraction]`` extra, so naming its types here (even under
    ``TYPE_CHECKING``) would fail mypy on an install without it.

    Indirected for test monkeypatching, like ``_try_import_docling``. ``None``
    covers both "docling is not installed" and "this docling build predates the
    factory" — the caller then leaves ``ocr_options`` unset and lets docling
    apply its own default, which is the tolerant behaviour the pre-#248 code had
    for missing option classes.
    """
    try:
        from docling.models.factories import get_ocr_factory

        return get_ocr_factory(allow_external_plugins=False)
    except Exception:
        return None


def _resolve_ocr_options(plan: OcrPlan, languages: list[str]) -> Any:
    """Resolve ``plan.engine_kind`` to a docling OCR options object.

    Returns ``None`` when OCR is disabled or the factory is unavailable — in
    both cases the caller omits ``ocr_options`` entirely rather than guessing.

    Raises ``ExtractorConfigurationError`` for a kind docling does not register.
    Validating against the *live* registry rather than a literal list in our
    config is what keeps this correct across docling upgrades — engines get
    added and renamed (this build knows ``tesserocr``, not ``tesseract_cli``).
    """
    if plan.engine_kind is None:
        return None
    factory = _try_import_ocr_factory()
    if factory is None:
        return None

    known = list(factory.registered_kind)
    if plan.engine_kind not in known:
        raise ExtractorConfigurationError(
            unknown_engine_message(plan.engine_kind, known)
        )
    return factory.create_options(kind=plan.engine_kind, lang=languages)


def _build_pdf_pipeline_options(cfg: SearchConfig) -> Any:
    """Build docling ``PdfPipelineOptions`` for ``cfg``, or ``None`` when the
    option classes are unavailable (older docling — caller falls back to a bare
    ``DocumentConverter()``).

    ``do_ocr`` and the engine both come from ``search.extractor_ocr_engine`` via
    the pure ``plan_ocr``. The pre-#248 code hardcoded ``do_ocr=True`` with
    ``EasyOcrOptions``, which overrode docling's own ``OcrAutoOptions`` default —
    the one that degrades to no-OCR instead of raising when no engine is
    installed.
    """
    try:
        from docling.datamodel.pipeline_options import PdfPipelineOptions
    except Exception:
        return None

    plan = plan_ocr(cfg.extractor_ocr_engine)
    ocr_options = _resolve_ocr_options(plan, cfg.extractor_ocr_languages)
    if ocr_options is None:
        return PdfPipelineOptions(do_ocr=plan.do_ocr)
    return PdfPipelineOptions(do_ocr=plan.do_ocr, ocr_options=ocr_options)


_TRANSIENT_THIRD_PARTY_MODULES: frozenset[str] = frozenset(
    {"requests", "httpx", "huggingface_hub", "urllib3", "aiohttp"}
)
"""Top-level package names whose exceptions signal a *transient* docling
failure (network blip, model-download hiccup) rather than a poison-pill blob.

These third-party classes are NOT in the builtin ``ConnectionError`` /
``TimeoutError`` hierarchy, so ``extract_worker._is_transient`` cannot
recognise them on its own. ``DoclingExtractor.extract`` opts them into
``TransientExtractorError`` here — keeping the wrapper-specific knowledge in
the wrapper and ``_TRANSIENT_EXC_TYPES`` narrow (its builtin guarantee).
Deliberately excludes broad packages (``os``, ``builtins``) that carry both
transient and permanent failures."""


def iter_exc_chain(exc: BaseException) -> Iterator[BaseException]:
    """Yield ``exc`` and each exception in its cause/context chain, in order.

    Follows ``__cause__`` first (explicit ``raise X from Y``) and falls back to
    ``__context__`` only when ``__suppress_context__`` is False — so
    ``raise X from None`` stops the walk, matching Python's own
    traceback-printing behaviour. A ``seen`` set guards against pathological
    reference cycles. Pure: no IO, reusable by any chain-inspecting caller
    (``_exc_chain_has_transient_module`` here, ``extract_worker._is_transient``).
    """
    seen: set[int] = set()
    cur: BaseException | None = exc
    while cur is not None and id(cur) not in seen:
        seen.add(id(cur))
        yield cur
        nxt = cur.__cause__
        if nxt is None and not cur.__suppress_context__:
            nxt = cur.__context__
        cur = nxt


def _exc_chain_has_transient_module(
    exc: BaseException,
    modules: frozenset[str] = _TRANSIENT_THIRD_PARTY_MODULES,
) -> bool:
    """True iff ``exc`` or any exception in its cause/context chain belongs to
    a top-level package in ``modules`` (``requests.exceptions`` → ``requests``).
    """
    return any(
        type(e).__module__.split(".", 1)[0] in modules
        for e in iter_exc_chain(exc)
    )


def _exc_chain_has_import_error(exc: BaseException) -> bool:
    """True iff ``exc`` or anything in its cause/context chain is an
    ``ImportError``.

    An ``ImportError`` escaping ``docling.convert()`` always means an optional
    docling dependency is absent — in practice an OCR engine, which docling
    constructs lazily inside the pipeline and which is not one of its own
    requirements. Matching on the *type* rather than on the message text
    (``"EasyOCR is not installed..."``) keeps this working across docling
    releases and across engines, each of which words its own message.
    """
    return any(isinstance(e, ImportError) for e in iter_exc_chain(exc))


class DoclingExtractor:
    """PDF-only extractor using docling for OCR + complex-PDF layout.

    Triggered by the extract_worker only when LightweightExtractor
    returned empty text or raised on a PDF. Lazy-imports docling so
    the package stays in the [extraction] optional dependency group.
    """

    name = "docling"
    version = "1.0"  # overwritten by importlib.metadata at extract time.

    def __init__(self, cfg: "SearchConfig | None" = None) -> None:
        """Construct with an optional SearchConfig.

        Defaults to SearchConfig() when omitted. The configured
        extractor_docling_max_pages and extractor_ocr_languages are read
        when .extract() builds the docling pipeline.
        """
        self._cfg = cfg if cfg is not None else SearchConfig()

    def supports(self, mime_type: str | None, filename: str | None) -> bool:
        """True iff the blob is a PDF (by MIME or by filename extension)."""
        return is_pdf(mime_type, filename)

    def extract(
        self, blob_path: Path, mime_type: str | None, *, filename: str | None = None
    ) -> ExtractedText:
        """Extract text from a PDF blob via docling.

        Raises ExtractorError when docling is not installed. Returns
        ExtractedText with text='' when docling produces no text (e.g.,
        the OCR pipeline fails to find any glyphs on a blank scan).

        Passes extractor_ocr_languages from SearchConfig into
        PdfPipelineOptions when the installed docling version exposes those
        option classes (falling back to a default DocumentConverter() on older
        builds — the language list is then best-effort). The page cap
        extractor_docling_max_pages is a convert-level argument on docling 2.x,
        so it is forwarded to converter.convert(..., max_num_pages=...) and
        applies on both the option-class and fallback paths.
        """
        DocumentConverter = _try_import_docling()
        if DocumentConverter is None:
            raise ExtractorError(
                "docling not installed; install via "
                "`uv sync --extra extraction`"
            )

        try:
            from importlib.metadata import PackageNotFoundError, version as pkg_version
            try:
                self_version = pkg_version("docling")
            except PackageNotFoundError:
                self_version = self.version
        except Exception:
            self_version = self.version

        # Build pipeline options if the docling option-classes are importable.
        # If the installed docling version doesn't expose these names, fall back
        # to default DocumentConverter() so older builds still work.
        #
        # ExtractorConfigurationError (an unknown engine name) propagates rather
        # than being swallowed by the fallback: it is the operator's typo, and
        # silently running with docling's default would hide it forever.
        pipeline_options = _build_pdf_pipeline_options(self._cfg)
        converter = None
        if pipeline_options is not None:
            try:
                from docling.datamodel.base_models import InputFormat
                from docling.document_converter import PdfFormatOption

                converter = DocumentConverter(
                    format_options={
                        InputFormat.PDF: PdfFormatOption(
                            pipeline_options=pipeline_options
                        ),
                    },
                )
            except Exception:
                converter = None
        if converter is None:
            converter = DocumentConverter()

        try:
            # The page cap is a convert-level argument on docling 2.x, not a
            # PdfPipelineOptions field — passing it here is what makes
            # extractor_docling_max_pages actually bound the OCR workload.
            result = converter.convert(
                str(blob_path),
                max_num_pages=self._cfg.extractor_docling_max_pages,
            )
        except Exception as exc:
            # A missing OCR engine surfaces as ImportError out of convert() —
            # docling constructs the engine lazily inside the pipeline. It is an
            # install problem, never the blob's fault, so it must not burn
            # retry_count (#248).
            if _exc_chain_has_import_error(exc):
                warn_ocr_engine_unavailable(
                    self._cfg.extractor_ocr_engine, str(exc)
                )
                raise ExtractorConfigurationError(
                    f"docling OCR engine unavailable: {exc}"
                ) from exc
            if _exc_chain_has_transient_module(exc):
                raise TransientExtractorError(
                    f"docling.convert transient failure: {exc}"
                ) from exc
            raise ExtractorError(f"docling.convert failed: {exc}") from exc

        try:
            text = result.document.export_to_markdown()
        except Exception as exc:
            raise ExtractorError(
                f"docling export_to_markdown failed: {exc}"
            ) from exc

        page_count = None
        if hasattr(result.document, "pages"):
            try:
                page_count = len(result.document.pages)
            except Exception:
                page_count = None

        return ExtractedText(
            text=(text or "").strip(),
            page_count=page_count,
            extractor=f"{self.name}@{self_version}",
        )
